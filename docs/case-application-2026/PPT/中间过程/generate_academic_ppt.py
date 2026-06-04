from __future__ import annotations

import subprocess
from pathlib import Path
from shutil import copy2

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
REPO_ROOT = ROOT.parents[2]
OUT = ROOT / "创AI-演示视频PPT-终稿.pptx"
BACKUP = ROOT / "中间过程" / "创AI-演示视频PPT-终稿.学术版生成前备份.pptx"
IMG = ROOT / "图资源"

NAVY = "0B2F5B"
BLUE = "1F5AA6"
PALE = "EAF2FB"
GREEN = "2E7D58"
GRAY = "5E6A75"
LIGHT = "D6E4F2"
TEXT = "1E2933"
WHITE = "FFFFFF"
FONT = "Microsoft YaHei"
SERIF = "SimSun"


def rgb(hex_color: str):
    from pptx.dml.color import RGBColor

    hex_color = hex_color.replace("#", "")
    return RGBColor(int(hex_color[:2], 16), int(hex_color[2:4], 16), int(hex_color[4:], 16))


def add_text(slide, text, x, y, w, h, size=14, color=TEXT, bold=False, align=None, font=FONT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    for para in tf.paragraphs:
        for r in para.runs:
            r.font.name = font
            r.font.size = Pt(size)
            r.font.color.rgb = rgb(color)
            r.font.bold = bold
    return box


def add_bullets(slide, items, x, y, w, h, size=13.2):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(TEXT)
        p.space_after = Pt(8)
    return box


def add_panel(slide, x, y, w, h, title, bullets):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(WHITE)
    shape.line.color.rgb = rgb(LIGHT)
    shape.line.width = Pt(1)
    add_text(slide, title, x + 0.18, y + 0.15, w - 0.36, 0.3, 12.5, BLUE, True)
    add_bullets(slide, bullets, x + 0.18, y + 0.6, w - 0.36, h - 0.75, 11.8)


def add_header(slide, section, title, page):
    add_text(slide, section, 0.45, 0.18, 1.25, 0.25, 8.5, GRAY)
    add_text(slide, title, 0.45, 0.42, 8.8, 0.45, 19, NAVY, True)
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.45), Inches(0.92), Inches(12.45), Inches(0.01))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb(LIGHT)
    line.line.color.rgb = rgb(LIGHT)
    add_text(slide, "XEdu Client · 创AI案例", 10.0, 0.2, 2.4, 0.25, 8.5, GRAY, align=PP_ALIGN.RIGHT)
    add_text(slide, f"{page:02d}", 12.35, 7.08, 0.4, 0.2, 8.5, GRAY, align=PP_ALIGN.RIGHT)


def add_caption(slide, text, x, y, w):
    add_text(slide, text, x, y, w, 0.25, 8.8, GRAY, align=PP_ALIGN.CENTER, font=SERIF)


def add_image_fit(slide, path, x, y, w, h):
    with Image.open(path) as img:
        iw, ih = img.size
    ratio = min(w / (iw / 96), h / (ih / 96))
    width = iw / 96 * ratio
    height = ih / 96 * ratio
    left = x + (w - width) / 2
    top = y + (h - height) / 2
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(width), Inches(height))


def add_figure_slide(prs, section, title, page, figure, caption, bullets, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, section, title, page)
    add_panel(slide, 0.55, 1.2, 3.55, 5.35, "报告要点", bullets)
    add_image_fit(slide, figure, 4.35, 1.2, 8.0, 5.35)
    add_caption(slide, caption, 4.35, 6.6, 8.0)
    if note:
        add_text(slide, note, 0.6, 6.72, 8.4, 0.22, 8.4, GRAY)


def add_two_image_slide(prs, section, title, page, left_figure, right_figure, left_caption, right_caption, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header(slide, section, title, page)
    add_panel(slide, 0.55, 1.2, 3.2, 5.35, "证据说明", bullets)
    add_image_fit(slide, left_figure, 4.05, 1.18, 4.25, 5.35)
    add_image_fit(slide, right_figure, 8.55, 1.18, 4.25, 5.35)
    add_caption(slide, left_caption, 4.05, 6.58, 4.25)
    add_caption(slide, right_caption, 8.55, 6.58, 4.25)


def build():
    if OUT.exists() and not BACKUP.exists():
        copy2(OUT, BACKUP)

    subprocess.run(["python3", str(ROOT / "中间过程" / "build_project_evidence_assets.py")], check=True, cwd=REPO_ROOT)

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    # 1 cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_text(slide, "XEdu Client", 0.7, 0.72, 6.2, 0.58, 28, NAVY, True)
    add_text(slide, "面向人工智能课堂的一体化实验学习平台", 0.72, 1.3, 7.0, 0.35, 16, TEXT)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.72), Inches(1.9), Inches(2.1), Inches(0.03))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(GREEN)
    add_text(slide, "创AI案例 · 人工智能学习工具 · 初中人工智能通识课", 0.72, 2.28, 6.7, 0.3, 12, GRAY)
    add_text(slide, "王海涛\n上海科技大学附属学校", 0.72, 5.75, 4.0, 0.55, 12, TEXT)
    add_image_fit(slide, IMG / "figure_system_framework_clear.png", 7.0, 0.7, 5.6, 5.6)
    add_caption(slide, "图示：平台总体框架概览", 7.0, 6.42, 5.6)

    add_figure_slide(
        prs,
        "01 研究背景",
        "课堂实验中的工具分散与路径断裂问题",
        2,
        IMG / "figure_1_project_background.png",
        "图1 课堂痛点与平台回应",
        [
            "人工智能课堂涉及讲解、积木、代码、素材和结果等多个入口。",
            "工具分散会放大环境配置成本，也会打断学生的任务上下文。",
            "本项目以统一实验入口回应课堂实施稳定性与资源复用问题。",
        ],
    )
    add_figure_slide(
        prs,
        "02 系统设计",
        "XEdu Client 的整体框架",
        3,
        IMG / "figure_system_framework_clear.png",
        "图2 课程资源、平台支撑层与课堂学习链路的关系",
        [
            "桌面端承接学生实验与教师备课两个入口。",
            "HTML、Blockly、Jupyter 形成递进式学习工作台。",
            "Flask 本地服务与 XEduHub 工具箱支撑课程资源运行。",
        ],
    )
    add_figure_slide(
        prs,
        "03 学习路径",
        "从讲解体验到代码实践的连续学习链路",
        4,
        IMG / "figure_3_learning_path.png",
        "图3 HTML、Blockly、Jupyter 的递进学习路径",
        [
            "HTML 讲解页用于建立任务情境和问题感知。",
            "Blockly 用可视化流程降低算法理解门槛。",
            "Jupyter 承接代码阅读、运行和结果复现。",
        ],
    )
    add_two_image_slide(
        prs,
        "04 课程资源",
        "真实课程包的结构化组织方式",
        5,
        IMG / "figure_course_page_evidence.png",
        IMG / "figure_course_resource_inventory.png",
        "图4 课程讲解页截图",
        "图5 课程包文件结构",
        [
            "课程样例来自项目真实目录，而不是另行制作的展示页。",
            "讲解页、Blockly、Notebook、Python、素材和输出共同构成一门课程。",
            "资源包结构使同一课程能够迁移、复用和继续修订。",
        ],
    )
    add_figure_slide(
        prs,
        "05 Blockly 证据",
        "可视化流程承接算法理解",
        6,
        IMG / "figure_blockly_workflow_evidence.png",
        "图6 从真实 Blockly XML 提取的任务流程",
        [
            "学生可先用积木看清“读取、检测、提框、分类、展示”的流程。",
            "Blockly 工作区与 Notebook 指向同一任务和同一素材。",
            "低门槛入口并不替代代码，而是为代码实践搭桥。",
        ],
    )
    add_two_image_slide(
        prs,
        "06 代码与运行",
        "Notebook/Python 与结果图片形成可复验证据",
        7,
        IMG / "figure_notebook_code_evidence.png",
        IMG / "figure_sample_output_evidence.png",
        "图7 Notebook 与 Python 脚本",
        "图8 输入、检测与裁剪结果",
        [
            "Notebook 保留课堂入口：导入、运行、查看 results。",
            "Python 脚本提供离线演示模式，方便录课和无模型环境复现。",
            "运行结果图来自课程目录，可作为视频演示与报告证据。",
        ],
    )
    add_figure_slide(
        prs,
        "07 课程共建",
        "多教师围绕同一课程包持续迭代",
        8,
        IMG / "figure_5b_collaborative_course.png",
        "图9 多教师课程共建与版本演进机制",
        [
            "共享课程仓库支持拉取、课堂使用、上传修订。",
            "课程从单次使用转为可复用、可迭代的教学资产。",
            "适合教研组围绕同一门课共同上课、共同改课。",
        ],
        "应用基础：已围绕人工智能课堂实验开展试用，服务 30 名学生，覆盖 12 节课。",
    )
    add_figure_slide(
        prs,
        "08 AI 辅助开发",
        "生成式人工智能在开发与材料整理中的作用",
        9,
        IMG / "figure_ai_development_process_clean.png",
        "图10 AI 辅助需求分析、测试清单和文稿优化",
        [
            "DeepSeek 用于辅助梳理课堂问题、测试清单与材料表达。",
            "AI 输出由教师团队筛选、修订并结合真实课堂目标落地。",
            "最终成果仍以可运行工具、真实课程资源和应用证据为准。",
        ],
    )
    add_figure_slide(
        prs,
        "09 结论",
        "应用价值与后续改进方向",
        10,
        IMG / "figure_application_value_clean.png",
        "图11 面向学生、教师和课程建设的应用价值",
        [
            "学生侧：降低实验入门门槛，增强学习连续性。",
            "教师侧：减少课堂组织成本，提高实施稳定性。",
            "课程侧：沉淀共享资源，支撑持续迭代与推广。",
        ],
        "后续将继续补充更多课程样例、课堂反馈和跨学段应用数据。",
    )

    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP

    prs.save(OUT)


if __name__ == "__main__":
    build()
