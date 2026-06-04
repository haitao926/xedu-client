from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
FIG = ROOT / "assets" / "figures"
OUT = ROOT / "deliverables" / "创AI-演示视频PPT-论文图版.pptx"

W, H = Inches(13.333), Inches(7.5)
NAVY = RGBColor(24, 52, 82)
TEAL = RGBColor(18, 112, 98)
MUTED = RGBColor(82, 100, 118)
LIGHT = RGBColor(244, 248, 250)
ORANGE = RGBColor(202, 98, 42)
WHITE = RGBColor(255, 255, 255)
PLACEHOLDER_BG = RGBColor(255, 244, 230)
PLACEHOLDER_LINE = RGBColor(224, 124, 52)


def set_font(run, size=20, bold=False, color=NAVY, name="Microsoft YaHei"):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_bg(slide, accent=TEAL):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    footer = slide.shapes.add_textbox(Inches(0.42), Inches(7.08), Inches(12.2), Inches(0.22))
    p = footer.text_frame.paragraphs[0]
    p.text = "创AI案例 · XEdu Client"
    p.alignment = PP_ALIGN.RIGHT
    set_font(p.runs[0], size=8, color=RGBColor(150, 160, 170))


def add_title(slide, title, subtitle=None, num=None):
    if num:
        n = slide.shapes.add_textbox(Inches(0.45), Inches(0.38), Inches(0.65), Inches(0.45))
        p = n.text_frame.paragraphs[0]
        p.text = num
        set_font(p.runs[0], size=20, bold=True, color=TEAL)
    box = slide.shapes.add_textbox(Inches(1.05), Inches(0.34), Inches(11.2), Inches(0.56))
    p = box.text_frame.paragraphs[0]
    p.text = title
    set_font(p.runs[0], size=28, bold=True, color=NAVY)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(1.08), Inches(0.92), Inches(10.8), Inches(0.35))
        p2 = sub.text_frame.paragraphs[0]
        p2.text = subtitle
        set_font(p2.runs[0], size=12, color=MUTED)


def add_bullets(slide, items, x, y, w, h, *, title=None):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = RGBColor(220, 229, 235)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.14)
    if title:
        p = tf.paragraphs[0]
        p.text = title
        set_font(p.runs[0], size=15, bold=True, color=TEAL)
    for idx, item in enumerate(items):
        p = tf.add_paragraph() if title or idx else tf.paragraphs[0]
        p.text = item
        p.level = 0
        p.space_after = Pt(7)
        set_font(p.runs[0], size=13, color=NAVY)


def add_image(slide, rel, x, y, w, h):
    p = FIG / rel
    slide.shapes.add_picture(str(p), x, y, width=w, height=h)


def add_label(slide, text, x, y, w, h, color=TEAL):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.fill.background()
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    set_font(p.runs[0], size=13, bold=True, color=WHITE)


def add_placeholder(slide, title, items, x, y, w, h):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = PLACEHOLDER_BG
    box.line.color.rgb = PLACEHOLDER_LINE
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.16)
    tf.margin_right = Inches(0.16)
    tf.margin_top = Inches(0.11)
    p = tf.paragraphs[0]
    p.text = f"【待补】{title}"
    set_font(p.runs[0], size=13, bold=True, color=ORANGE)
    for item in items:
        para = tf.add_paragraph()
        para.text = item
        para.space_after = Pt(3)
        set_font(para.runs[0], size=10.5, color=NAVY)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    # 1 Cover
    s = blank(prs)
    add_bg(s, NAVY)
    title = s.shapes.add_textbox(Inches(0.9), Inches(1.45), Inches(11.8), Inches(1.2))
    p = title.text_frame.paragraphs[0]
    p.text = "XEdu Client"
    set_font(p.runs[0], size=50, bold=True, color=NAVY)
    sub = s.shapes.add_textbox(Inches(0.95), Inches(2.55), Inches(10.8), Inches(0.8))
    p = sub.text_frame.paragraphs[0]
    p.text = "面向课堂实验的 Jupyter + Blockly 一体化人工智能学习工具"
    set_font(p.runs[0], size=23, color=TEAL)
    meta = s.shapes.add_textbox(Inches(0.98), Inches(3.45), Inches(8.6), Inches(0.55))
    p = meta.text_frame.paragraphs[0]
    p.text = "创AI案例 · 人工智能学习工具 · 初中人工智能通识课"
    set_font(p.runs[0], size=15, color=MUTED)
    add_placeholder(
        s,
        "封面信息",
        ["作者：[待填写]", "单位：[待填写]", "应用学校/班级：[待填写]"],
        Inches(0.98),
        Inches(3.95),
        Inches(4.45),
        Inches(0.9),
    )
    add_image(s, "paper-roil/figure_2_architecture.png", Inches(6.85), Inches(1.18), Inches(5.75), Inches(3.24))
    add_label(s, "实验环境承载", Inches(0.95), Inches(5.05), Inches(2.0), Inches(0.5), TEAL)
    add_label(s, "学习主线", Inches(3.2), Inches(5.05), Inches(2.0), Inches(0.5), NAVY)
    add_label(s, "XEduHub工具箱", Inches(5.45), Inches(5.05), Inches(2.2), Inches(0.5), TEAL)
    add_label(s, "课程共建更新", Inches(7.95), Inches(5.05), Inches(2.2), Inches(0.5), ORANGE)

    # 2 Problem
    s = blank(prs)
    add_bg(s)
    add_title(s, "为什么需要平台化实验环境", "四类课堂问题对应四类平台能力", "01")
    add_image(s, "paper-roil/figure_1_problem_solution.png", Inches(0.75), Inches(1.35), Inches(7.55), Inches(4.25))
    add_bullets(s, ["环境、工具、资源分散，课堂进入成本高", "讲解、积木、代码之间缺少连续路径", "课程需要共享、复用、持续更新"], Inches(8.55), Inches(1.75), Inches(3.95), Inches(2.1), title="核心判断")
    add_bullets(s, ["不是单点工具", "而是课堂实验平台"], Inches(8.55), Inches(4.2), Inches(3.95), Inches(1.15), title="案例定位")

    # 3 Architecture
    s = blank(prs)
    add_bg(s)
    add_title(s, "总体架构：桌面端 + 本地服务 + 课程资源", "用统一工作台承载人工智能实验环境", "02")
    add_image(s, "paper-roil/figure_2_architecture.png", Inches(0.72), Inches(1.3), Inches(8.15), Inches(4.58))
    add_bullets(s, ["Electron 提供桌面端入口", "HTML / Blockly / Jupyter 形成教学工作台", "XEduHub、课程目录、云端仓库支撑平台化扩展"], Inches(9.05), Inches(1.6), Inches(3.55), Inches(2.7), title="架构要点")

    # 4 Learning path
    s = blank(prs)
    add_bg(s)
    add_title(s, "学习主线：HTML—Blockly—Jupyter", "不是三个入口并列，而是体验、实践、应用创新的递进路径", "03")
    add_image(s, "paper-roil/figure_3_learning_path.png", Inches(0.85), Inches(1.45), Inches(8.5), Inches(4.78))
    add_bullets(s, ["HTML：任务情境与问题感知", "Blockly：低门槛理解检测—裁剪—分类流程", "Jupyter：阅读代码、修改参数、复现实验"], Inches(9.5), Inches(1.7), Inches(3.05), Inches(3.05), title="课堂价值")

    # 5 XEduHub
    s = blank(prs)
    add_bg(s)
    add_title(s, "XEduHub：人工智能实验能力底座", "把输入、模型、流程和结果从零散脚本中抽象出来", "04")
    add_image(s, "paper-roil/figure_4_xeduhub.png", Inches(0.85), Inches(1.35), Inches(8.0), Inches(4.5))
    add_bullets(s, ["同时支撑 Blockly 与 Jupyter", "沉淀可复用 AI 实验能力", "便于后续扩展更多课程样例"], Inches(9.1), Inches(1.75), Inches(3.45), Inches(2.45), title="工具箱作用")

    # 6 Course update
    s = blank(prs)
    add_bg(s)
    add_title(s, "课程资源：从个人文件夹到可迭代资产", "上一节、改一节、同步一节", "05")
    add_image(s, "paper-roil/figure_5_course_update.png", Inches(0.8), Inches(1.35), Inches(5.75), Inches(3.24))
    add_image(s, "paper-roil/figure_5b_collaborative_course.png", Inches(6.85), Inches(1.35), Inches(5.75), Inches(3.24))
    add_bullets(s, ["支持多教师围绕同一门课共同维护", "课程可拉取、修订、上传、形成版本演进", "教研组可共享一门课、共上一门课、共改一门课"], Inches(1.05), Inches(5.18), Inches(11.4), Inches(1.3), title="协同机制")
    add_placeholder(
        s,
        "真实共建证据",
        ["共同维护教师人数：[待填写]", "课程版本/更新记录截图：[待补]", "共享课程仓库或课程列表截图：[待补]"],
        Inches(8.35),
        Inches(5.12),
        Inches(4.0),
        Inches(1.48),
    )

    # 7 Classroom flow
    s = blank(prs)
    add_bg(s)
    add_title(s, "典型课堂应用流程", "教师侧与学生侧在同一平台完成课前、课中、课后闭环", "06")
    add_image(s, "paper-roil/figure_6_classroom_flow.png", Inches(0.85), Inches(1.32), Inches(8.2), Inches(4.62))
    add_bullets(s, ["课前：教师组织课程资源", "课中：学生进入实验，完成 Blockly 与 Jupyter 实践", "课后：教师更新课程，沉淀后续班级可复用资源"], Inches(9.3), Inches(1.75), Inches(3.2), Inches(2.8), title="实施路径")
    add_placeholder(
        s,
        "课堂应用数据",
        ["应用学校：[待填写]", "班级/人数：[待填写]", "实施课时：[待填写]"],
        Inches(9.3),
        Inches(4.75),
        Inches(3.2),
        Inches(1.2),
    )

    # 8 Sample
    s = blank(prs)
    add_bg(s)
    add_title(s, "样例任务：目标检测与细粒度分类", "用真实课程素材展示输入、检测、裁剪和代码实践链条", "07")
    add_image(s, "figure_7_sample_experiment.png", Inches(0.75), Inches(1.25), Inches(8.45), Inches(4.75))
    add_bullets(s, ["演示主线固定为 01_目标检测与细粒度分类", "从 Blockly 流程理解过渡到 Jupyter 代码复现", "后续可补课堂照片、学生作品、AI 辅助开发证据"], Inches(9.45), Inches(1.7), Inches(3.0), Inches(3.1), title="证据补充")
    add_placeholder(
        s,
        "真实素材替换",
        ["Blockly 截图：[待补]", "Jupyter 运行截图：[待补]", "学生作品/结果图：[待补]"],
        Inches(9.45),
        Inches(4.95),
        Inches(3.0),
        Inches(1.25),
    )

    # 9 Effect and innovation
    s = blank(prs)
    add_bg(s, ORANGE)
    add_title(s, "应用成效与创新总结", "降低门槛、增强连续性、提升课程复现与共建能力", "08")
    add_bullets(s, ["学生侧：先理解流程，再进入代码，降低人工智能实验起步门槛", "教师侧：减少环境切换与文件寻找，把课堂时间还给任务探究", "课程侧：资源可共享、可修订、可同步，支持教研组共同维护同一门课"], Inches(0.95), Inches(1.55), Inches(5.5), Inches(2.8), title="应用成效")
    add_bullets(s, ["创新一：HTML—Blockly—Jupyter 连续学习主线", "创新二：XEduHub 统一实验能力底座", "创新三：面向教师的课程分享与共同维护机制"], Inches(6.85), Inches(1.55), Inches(5.55), Inches(2.8), title="核心创新")
    add_bullets(s, ["待补：应用学校、班级数量、学生人数、课时数", "待补：国产 AI 工具名称、提示词、对话截图", "待补：课堂照片、学生作品、反馈摘录"], Inches(0.95), Inches(4.85), Inches(11.45), Inches(1.3), title="提交前证据位")
    add_placeholder(
        s,
        "必须补齐后再提交",
        ["课堂照片：避免学生正脸", "国产 AI 工具与提示词截图", "真实反馈或学生作品"],
        Inches(7.2),
        Inches(4.93),
        Inches(5.0),
        Inches(1.12),
    )

    # 10 Closing
    s = blank(prs)
    add_bg(s, NAVY)
    box = s.shapes.add_textbox(Inches(1.05), Inches(2.2), Inches(11.2), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    p.text = "谢谢观看"
    p.alignment = PP_ALIGN.CENTER
    set_font(p.runs[0], size=46, bold=True, color=NAVY)
    sub = s.shapes.add_textbox(Inches(1.4), Inches(3.2), Inches(10.5), Inches(0.75))
    p = sub.text_frame.paragraphs[0]
    p.text = "XEdu Client：面向人工智能课堂的一体化实验学习平台"
    p.alignment = PP_ALIGN.CENTER
    set_font(p.runs[0], size=19, color=TEAL)
    add_label(s, "创AI案例 · 人工智能学习工具", Inches(5.0), Inches(4.45), Inches(3.35), Inches(0.5), TEAL)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
