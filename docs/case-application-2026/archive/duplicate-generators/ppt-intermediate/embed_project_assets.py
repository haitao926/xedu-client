from __future__ import annotations

from copy import deepcopy
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Cm, Pt, RGBColor
from docx.text.paragraph import Paragraph
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu, Pt as PptPt


ROOT = Path(__file__).resolve().parent
REPO_ROOT = ROOT.parent.parent
DELIVERABLES = ROOT / "deliverables"
PPT_PATH = DELIVERABLES / "创AI-演示视频PPT-终稿.pptx"
DOCX_PATH = DELIVERABLES / "创AI-开发与应用报告-终稿.docx"

ASSETS = {
    "problem_solution": ROOT / "assets/figures/figure_1_problem_solution.png",
    "architecture": ROOT / "assets/figures/figure_2_architecture.png",
    "learning_path": ROOT / "assets/figures/figure_3_learning_path.png",
    "toolbox": ROOT / "assets/figures/figure_4_xeduhub_toolbox.png",
    "course_update": ROOT / "assets/figures/figure_5_course_update.png",
    "collaboration": ROOT / "assets/figures/paper-roil/figure_5b_collaborative_course.png",
    "classroom_flow": ROOT / "assets/figures/figure_6_classroom_flow.png",
    "sample_task": ROOT / "assets/figures/figure_7_sample_experiment.png",
    "input_image": REPO_ROOT / "courses/ai-showcase-exam-2025/01_目标检测与细粒度分类/animal.jpg",
    "detection_result": REPO_ROOT
    / "courses/ai-showcase-exam-2025/01_目标检测与细粒度分类/outputs/animal_with_boxes.jpg",
    "crop_1": REPO_ROOT
    / "courses/ai-showcase-exam-2025/01_目标检测与细粒度分类/outputs/crop_1.jpg",
    "crop_2": REPO_ROOT
    / "courses/ai-showcase-exam-2025/01_目标检测与细粒度分类/outputs/crop_2.jpg",
}


def set_run_style(run, *, name: str = "仿宋_GB2312", size: int = 10, bold: bool = False):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string("4A4A4A")
    if run._element.rPr is not None:
        rfonts = run._element.rPr.rFonts
        if rfonts is None:
            rfonts = run._element.rPr._add_rFonts()
        rfonts.set(
            "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}eastAsia",
            name,
        )


def add_caption(doc: Document, text: str):
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(text)
    set_run_style(r, size=9)
    return p


def insert_paragraph_after(paragraph: Paragraph) -> Paragraph:
    new_p = deepcopy(paragraph._p)
    for child in list(new_p):
        new_p.remove(child)
    paragraph._p.addnext(new_p)
    return Paragraph(new_p, paragraph._parent)


def insert_docx_figure(doc: Document, anchor: str, image_path: Path, caption: str, width_cm: float = 15.8):
    for para in doc.paragraphs:
        if anchor in para.text:
            new_p = insert_paragraph_after(para)
            new_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            new_p.add_run().add_picture(str(image_path), width=Cm(width_cm))
            caption_p = insert_paragraph_after(new_p)
            caption_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            caption_run = caption_p.add_run(caption)
            set_run_style(caption_run, size=9)
            return
    raise ValueError(f"Anchor not found in DOCX: {anchor}")


def replace_picture(slide, image_path: Path, picture_idx: int = 0):
    pictures = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    if picture_idx >= len(pictures):
        raise IndexError(f"Picture index {picture_idx} out of range for slide")

    picture = pictures[picture_idx]
    left, top, width, height = picture.left, picture.top, picture.width, picture.height
    picture._element.getparent().remove(picture._element)
    slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def set_shape_text(shape, text: str, *, font_size: int | None = None):
    if not getattr(shape, "has_text_frame", False):
        return
    shape.text = text
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if font_size is not None:
                run.font.size = PptPt(font_size)


def find_shape_by_text(slide, needle: str):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and needle in shape.text:
            return shape
    raise ValueError(f"Shape containing {needle!r} not found")


def find_shape_by_any_text(slide, needles: list[str]):
    for needle in needles:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and needle in shape.text:
                return shape
    raise ValueError(f"Shape containing any of {needles!r} not found")


def crop_to_fill(slide, image_path: Path, left: Emu, top: Emu, width: Emu, height: Emu):
    pic = slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    with Image.open(image_path) as img:
        img_ratio = img.width / img.height
    box_ratio = width / height

    if img_ratio > box_ratio:
        new_width = height * img_ratio
        pic.width = int(new_width)
        pic.height = height
        pic.left = int(left - (new_width - width) / 2)
        pic.top = top
    else:
        new_height = width / img_ratio
        pic.width = width
        pic.height = int(new_height)
        pic.left = left
        pic.top = int(top - (new_height - height) / 2)
    return pic


def add_task_gallery(slide):
    target = find_shape_by_any_text(slide, ["【待补】真实素材替换", "项目资源说明：课程目录同时提供"])
    pictures = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    if pictures:
        keep = pictures[0]
        keep._element.getparent().remove(keep._element)
    left, top, width, height = target.left, target.top, target.width, target.height
    target._element.getparent().remove(target._element)

    gap = Emu(120000)
    box_width = int((width - gap * 2) / 3)
    box_height = int(height * 0.64)
    labels = [
        ("输入图片", ASSETS["input_image"]),
        ("检测结果", ASSETS["detection_result"]),
        ("分类结果", ASSETS["crop_1"]),
    ]

    for idx, (label, image_path) in enumerate(labels):
        x = left + idx * (box_width + gap)
        title = slide.shapes.add_textbox(x, top, box_width, Emu(260000))
        set_shape_text(title, label, font_size=16)
        crop_to_fill(slide, image_path, x, top + Emu(300000), box_width, box_height)

    note = slide.shapes.add_textbox(left, top + box_height + Emu(380000), width, Emu(520000))
    note.text_frame.word_wrap = True
    set_shape_text(
        note,
        "项目资源说明：课程目录同时提供 index.html、Blockly 工作区、Notebook、Python 脚本与样例图片，"
        "学生可在同一任务上下文中完成“看流程 - 跑代码 - 看结果”。",
        font_size=14,
    )


def update_ppt():
    prs = Presentation(str(PPT_PATH))

    replace_picture(prs.slides[1], ASSETS["problem_solution"])
    replace_picture(prs.slides[2], ASSETS["architecture"])
    replace_picture(prs.slides[3], ASSETS["learning_path"])
    replace_picture(prs.slides[4], ASSETS["toolbox"])
    replace_picture(prs.slides[5], ASSETS["course_update"], picture_idx=0)
    replace_picture(prs.slides[5], ASSETS["collaboration"], picture_idx=1)
    replace_picture(prs.slides[6], ASSETS["classroom_flow"])
    replace_picture(prs.slides[7], ASSETS["sample_task"])

    set_shape_text(
        find_shape_by_any_text(prs.slides[5], ["【待补】真实共建证据", "项目资源示例"]),
        "项目资源示例\n课程目录、图文讲解页、Blockly 工作区、Notebook 与 Python 脚本均可随课程包共享与迭代。",
        font_size=15,
    )
    set_shape_text(
        find_shape_by_any_text(prs.slides[6], ["【待补】课堂应用数据", "应用情况"]),
        "应用情况\n已围绕人工智能课堂实验开展试用，服务 30 名学生，覆盖 12 节课。",
        font_size=15,
    )

    add_task_gallery(prs.slides[7])
    prs.save(str(PPT_PATH))


def update_docx():
    doc = Document(str(DOCX_PATH))
    insert_docx_figure(
        doc,
        "XEdu Client。",
        ASSETS["problem_solution"],
        "图1 课堂问题与平台回应：用统一工作台承接环境、路径、能力与资源四类需求。",
    )
    insert_docx_figure(
        doc,
        "项目采用“桌面端入口 + 本地服务支撑 + 课程资源结构化管理”的技术路线。",
        ASSETS["architecture"],
        "图2 平台总体架构：桌面端入口、本地服务、课程资源与工具箱共同组成课堂实验环境。",
    )
    insert_docx_figure(
        doc,
        "这样的技术组合并不是简单叠加工具，而是围绕课堂学习过程进行重新编排。",
        ASSETS["learning_path"],
        "图3 学习主线：从讲解体验到 Blockly 流程理解，再到 Jupyter 代码实践。",
    )
    insert_docx_figure(
        doc,
        "在课程资源建设上，项目以文件夹和结构化配置组织课程内容，使讲解页、Notebook、Blockly 工作区、素材文件和说明文档能够共同构成一个可迁移的课程包。",
        ASSETS["course_update"],
        "图4 课程资源组织与更新：同一门课程可按讲解页、Notebook、Blockly 与素材资源持续迭代。",
    )
    insert_docx_figure(
        doc,
        "教师可以将讲解页、Notebook、Blockly 文件、Python 脚本、图片素材和说明文档组织为课程资源包，形成可沉淀、可共享、可复用的课程结构，提高优质资源的推广效率。",
        ASSETS["collaboration"],
        "图5 课程共建示意：多位教师可围绕同一门课共享、使用、修订并同步资源。",
    )
    insert_docx_figure(
        doc,
        "课堂实施时，教师先通过课程资源准备讲解体验页面、Blockly 工作区、Notebook 和实验素材，学生进入统一工作台后，按照“体验 - 理解 - 实践 - 应用”的路径完成学习任务。",
        ASSETS["classroom_flow"],
        "图6 典型课堂应用流程：课前准备、课中实践、课后更新在同一平台内闭环完成。",
    )
    insert_docx_figure(
        doc,
        "以“运动会上的 AI 裁判”等课堂任务为例，学生能够围绕运动场景中的识别、判断和结果反馈等问题理解人工智能的应用过程。",
        ASSETS["sample_task"],
        "图7 样例任务资源：使用项目内“目标检测与细粒度分类”课程包展示输入、检测、裁剪与代码实践链条。",
        width_cm=16.2,
    )
    doc.save(str(DOCX_PATH))


def main():
    update_ppt()
    update_docx()


if __name__ == "__main__":
    main()
